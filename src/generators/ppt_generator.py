"""
PPTGenerator - PPT文件生成器
支持多种模板和样式
"""
import os
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class PPTTemplate:
    """PPT模板配置"""

    # 商务专业模板
    BUSINESS = {
        "name": "商务专业",
        "slide_width": Inches(10),
        "slide_height": Inches(5.625),
        "colors": {
            "primary": RGBColor(31, 78, 120),      # 深蓝
            "secondary": RGBColor(68, 114, 196),   # 蓝色
            "accent": RGBColor(237, 125, 49),      # 橙色
            "text": RGBColor(0, 0, 0),             # 黑色
            "background": RGBColor(255, 255, 255)  # 白色
        },
        "fonts": {
            "title": "微软雅黑",
            "content": "微软雅黑",
            "title_size": Pt(44),
            "subtitle_size": Pt(24),
            "content_title_size": Pt(32),
            "content_size": Pt(18)
        }
    }

    # 创意风格模板
    CREATIVE = {
        "name": "创意风格",
        "slide_width": Inches(10),
        "slide_height": Inches(5.625),
        "colors": {
            "primary": RGBColor(230, 57, 70),      # 红色
            "secondary": RGBColor(241, 250, 238),  # 浅绿
            "accent": RGBColor(255, 211, 105),     # 黄色
            "text": RGBColor(41, 50, 65),          # 深灰
            "background": RGBColor(255, 255, 255)
        },
        "fonts": {
            "title": "微软雅黑",
            "content": "微软雅黑",
            "title_size": Pt(48),
            "subtitle_size": Pt(26),
            "content_title_size": Pt(36),
            "content_size": Pt(20)
        }
    }

    # 学术风格模板
    ACADEMIC = {
        "name": "学术风格",
        "slide_width": Inches(10),
        "slide_height": Inches(5.625),
        "colors": {
            "primary": RGBColor(44, 62, 80),       # 深灰蓝
            "secondary": RGBColor(52, 73, 94),     # 灰蓝
            "accent": RGBColor(192, 57, 43),       # 深红
            "text": RGBColor(0, 0, 0),
            "background": RGBColor(255, 255, 255)
        },
        "fonts": {
            "title": "Arial",
            "content": "Arial",
            "title_size": Pt(40),
            "subtitle_size": Pt(22),
            "content_title_size": Pt(30),
            "content_size": Pt(16)
        }
    }


class PPTGenerator:
    """PPT生成器"""

    def __init__(self, template: str = "business"):
        """
        初始化PPT生成器

        Args:
            template: 模板名称(business/creative/academic)
        """
        self.prs = Presentation()

        # 选择模板
        if template.lower() == "creative":
            self.template = PPTTemplate.CREATIVE
        elif template.lower() == "academic":
            self.template = PPTTemplate.ACADEMIC
        else:
            self.template = PPTTemplate.BUSINESS

        # 设置页面尺寸
        self.prs.slide_width = self.template["slide_width"]
        self.prs.slide_height = self.template["slide_height"]

        print(f"📄 使用模板: {self.template['name']}")

    def create_presentation(
        self,
        outline: Dict,
        contents: List[Dict],
        images: List[str] = None
    ) -> str:
        """
        创建完整的PPT

        Args:
            outline: 大纲(包含title, slides等)
            contents: 各页内容列表
            images: 各页图片路径列表(可选)

        Returns:
            保存的文件路径
        """
        print(f"\n📊 开始创建PPT: {outline['title']}")

        for i, slide_info in enumerate(outline["slides"]):
            if i >= len(contents):
                print(f"   ⚠️  警告: 内容不足,跳过第{i+1}页")
                break

            content = contents[i]
            slide_type = content.get("type", "content")

            # 获取对应的图片路径
            image_path = None
            if images and i < len(images):
                image_path = images[i]

            # 根据类型添加幻灯片
            if slide_type == "cover":
                self.add_cover_slide(content)
            elif slide_type == "conclusion":
                self.add_conclusion_slide(content)
            else:
                self.add_content_slide(content, image_path)

            print(f"   ✅ 第{i+1}页: {content.get('title', '')}")

        # 保存文件
        filename = self._sanitize_filename(outline["title"])
        filepath = f"output/{filename}.pptx"
        self.prs.save(filepath)

        print(f"\n✅ PPT创建成功: {filepath}")
        return filepath

    def add_cover_slide(self, content: Dict) -> None:
        """
        添加封面页

        Args:
            content: 包含title, subtitle
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局

        colors = self.template["colors"]
        fonts = self.template["fonts"]

        # 背景色块(装饰)
        left = Inches(0)
        top = Inches(0)
        width = self.prs.slide_width
        height = Inches(2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.2), Inches(8), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.font.name = fonts["title"]
        title_para.font.size = fonts["title_size"]
        title_para.font.bold = True
        title_para.font.color.rgb = colors["primary"]
        title_para.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.6), Inches(8), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = content.get("subtitle", "")
        subtitle_frame.word_wrap = True

        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = fonts["content"]
        subtitle_para.font.size = fonts["subtitle_size"]
        subtitle_para.font.color.rgb = colors["text"]
        subtitle_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, content: Dict, image_path: str = None) -> None:
        """
        添加内容页

        Args:
            content: 包含title, content(要点列表)
            image_path: 图片路径(可选)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        colors = self.template["colors"]
        fonts = self.template["fonts"]

        # 标题区域
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.font.name = fonts["title"]
        title_para.font.size = fonts["content_title_size"]
        title_para.font.bold = True
        title_para.font.color.rgb = colors["primary"]

        # 标题下划线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.15),
            Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors["accent"]
        line.line.fill.background()

        # 判断是否有图片
        if image_path and os.path.exists(image_path):
            # 有图片:左侧内容,右侧图片
            self._add_content_with_image(slide, content, image_path, colors, fonts)
        else:
            # 无图片:全宽内容
            self._add_content_only(slide, content, colors, fonts)

    def _add_content_only(self, slide, content: Dict, colors: Dict, fonts: Dict) -> None:
        """添加纯文本内容"""
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.6), Inches(8), Inches(3.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = 1

        points = content.get("content", [])
        for i, point in enumerate(points):
            if i > 0:
                p = content_frame.add_paragraph()
            else:
                p = content_frame.paragraphs[0]

            p.text = f"• {point}"
            p.font.name = fonts["content"]
            p.font.size = fonts["content_size"]
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(12) if i > 0 else Pt(0)
            p.line_spacing = 1.3

    def _add_content_with_image(
        self,
        slide,
        content: Dict,
        image_path: str,
        colors: Dict,
        fonts: Dict
    ) -> None:
        """添加图文混排内容"""
        # 左侧文本区域(更窄)
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(4.5), Inches(3.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = 1

        points = content.get("content", [])
        for i, point in enumerate(points):
            if i > 0:
                p = content_frame.add_paragraph()
            else:
                p = content_frame.paragraphs[0]

            p.text = f"• {point}"
            p.font.name = fonts["content"]
            p.font.size = Pt(16)  # 稍小一点
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(10) if i > 0 else Pt(0)
            p.line_spacing = 1.2

        # 右侧图片
        try:
            from PIL import Image

            # 获取图片尺寸
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height

            # 计算合适的显示尺寸
            max_width = Inches(4.5)
            max_height = Inches(3.5)

            if aspect_ratio > 1:  # 横图
                pic_width = max_width
                pic_height = pic_width / aspect_ratio
            else:  # 竖图
                pic_height = max_height
                pic_width = pic_height * aspect_ratio

            # 居中对齐
            left = Inches(5.5) + (max_width - pic_width) / 2
            top = Inches(1.6) + (max_height - pic_height) / 2

            # 插入图片
            slide.shapes.add_picture(
                image_path,
                left, top,
                width=pic_width,
                height=pic_height
            )

        except Exception as e:
            print(f"      ⚠️  图片插入失败: {str(e)}")

    def add_conclusion_slide(self, content: Dict) -> None:
        """
        添加结束页

        Args:
            content: 包含title, content
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        colors = self.template["colors"]
        fonts = self.template["fonts"]

        # 背景装饰
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(1.8), Inches(6), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "谢谢")

        title_para = title_frame.paragraphs[0]
        title_para.font.name = fonts["title"]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # 副文本
        content_items = content.get("content", [])
        if content_items:
            content_box = slide.shapes.add_textbox(
                Inches(2), Inches(3.2), Inches(6), Inches(1.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, item in enumerate(content_items):
                if i > 0:
                    p = content_frame.add_paragraph()
                else:
                    p = content_frame.paragraphs[0]

                p.text = item
                p.font.name = fonts["content"]
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
                p.space_before = Pt(8) if i > 0 else Pt(0)

    def add_custom_slide(
        self,
        title: str,
        content: List[str],
        layout: str = "content"
    ) -> None:
        """
        添加自定义幻灯片

        Args:
            title: 标题
            content: 内容列表
            layout: 布局类型
        """
        content_dict = {
            "title": title,
            "content": content,
            "type": layout
        }

        if layout == "cover":
            self.add_cover_slide(content_dict)
        elif layout == "conclusion":
            self.add_conclusion_slide(content_dict)
        else:
            self.add_content_slide(content_dict)

    def _sanitize_filename(self, filename: str) -> str:
        """
        清理文件名,移除非法字符

        Args:
            filename: 原始文件名

        Returns:
            清理后的文件名
        """
        # 移除非法字符
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '')

        # 限制长度
        if len(filename) > 50:
            filename = filename[:50]

        return filename.strip()

    def get_slide_count(self) -> int:
        """获取当前幻灯片数量"""
        return len(self.prs.slides)

    def save(self, filepath: str) -> None:
        """
        保存PPT文件

        Args:
            filepath: 保存路径
        """
        self.prs.save(filepath)
        print(f"✅ PPT已保存: {filepath}")