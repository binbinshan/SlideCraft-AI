"""
SlideCraft AI - 第1周原型 (OpenAI版本)
完整的端到端演示: 主题 → GPT生成内容 → 创建PPT

运行: python prototype_demo.py
"""
import os
import json
import re
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 加载环境变量
load_dotenv()


class SlideCraftPrototype:
    """SlideCraft AI 原型类"""

    def __init__(self):
        """初始化"""
        api_key = os.getenv("DEEPSEEK_API_KEY")
        if not api_key:
            raise ValueError("请在.env文件中设置DEEPSEEK_API_KEY")

        self.client = OpenAI(api_key=api_key,base_url="https://api.deepseek.com")
        self.model = os.getenv("DEEPSEEK_MODEL", "deepseek-chat")
        self.presentation = None

    def generate_outline(self, topic, num_slides=5):
        """生成PPT大纲"""
        print(f"\n🤖 正在为主题'{topic}'生成大纲...")

        prompt = f"""
请为以下主题生成一个{num_slides}页的PPT大纲:

主题: {topic}

要求:
1. 第1页必须是封面(cover)
2. 最后1页必须是结束页(conclusion)
3. 中间是内容页(content)
4. 每页要有清晰的标题

请以JSON格式返回,只返回JSON,不要markdown标记:
{{
  "title": "PPT总标题",
  "slides": [
    {{"page": 1, "title": "标题", "type": "cover"}},
    {{"page": 2, "title": "标题", "type": "content"}},
    ...
  ]
}}
"""

        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "你是一个专业的PPT大纲设计师。请始终返回有效的JSON格式。"},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2048,
                temperature=0.7
            )

            response_text = response.choices[0].message.content.strip()
            # 去除可能的markdown标记
            response_text = re.sub(r'^```json\s*|\s*```$', '', response_text, flags=re.MULTILINE)

            outline = json.loads(response_text)
            print(f"✅ 大纲生成成功! 共{len(outline['slides'])}页")
            return outline

        except Exception as e:
            print(f"❌ 大纲生成失败: {str(e)}")
            raise

    def generate_slide_content(self, slide_info, full_topic):
        """为单页生成详细内容"""
        slide_type = slide_info.get("type", "content")
        title = slide_info.get("title", "")

        print(f"   生成第{slide_info['page']}页: {title}")

        if slide_type == "cover":
            # 封面页只需要标题和副标题
            return {
                "title": title,
                "subtitle": f"关于{full_topic}的深入探讨",
                "content": []
            }
        elif slide_type == "conclusion":
            # 结束页
            return {
                "title": title,
                "content": ["感谢观看!", "欢迎提问与讨论"]
            }
        else:
            # 内容页
            prompt = f"""
为PPT的这一页生成详细内容:

页面标题: {title}
整体主题: {full_topic}

要求:
1. 生成3-5个要点
2. 每个要点简洁明了,适合PPT展示
3. 每个要点控制在20字以内

返回JSON格式,只返回JSON:
{{
  "title": "{title}",
  "content": ["要点1", "要点2", "要点3"]
}}
"""

            try:
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": "你是一个专业的PPT内容创作者。请始终返回有效的JSON格式。"},
                        {"role": "user", "content": prompt}
                    ],
                    max_tokens=1024,
                    temperature=0.7
                )

                response_text = response.choices[0].message.content.strip()
                response_text = re.sub(r'^```json\s*|\s*```$', '', response_text, flags=re.MULTILINE)

                content = json.loads(response_text)
                return content

            except Exception as e:
                print(f"      ⚠️  生成失败,使用默认内容: {str(e)}")
                return {
                    "title": title,
                    "content": ["内容生成中...", "请稍候..."]
                }

    def create_ppt(self, outline, contents):
        """创建PPT文件"""
        print(f"\n📊 开始创建PPT...")

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        for i, slide_info in enumerate(outline["slides"]):
            content = contents[i]
            slide_type = slide_info.get("type", "content")

            # 添加空白幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            if slide_type == "cover":
                # 封面页
                self._add_cover_slide(slide, content)
            elif slide_type == "conclusion":
                # 结束页
                self._add_conclusion_slide(slide, content)
            else:
                # 内容页
                self._add_content_slide(slide, content)

        # 保存
        os.makedirs('output', exist_ok=True)
        output_path = f'output/{outline["title"]}.pptx'
        prs.save(output_path)
        print(f"✅ PPT创建成功! 保存路径: {output_path}")
        return output_path

    def _add_cover_slide(self, slide, content):
        """添加封面页"""
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = content["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.2), Inches(8), Inches(0.6)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = content.get("subtitle", "")
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.alignment = PP_ALIGN.CENTER

    def _add_conclusion_slide(self, slide, content):
        """添加结束页"""
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(2), Inches(6), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = content["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, slide, content):
        """添加内容页"""
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = content["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(3.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for point in content.get("content", []):
            p = content_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.space_before = Pt(12)

    def generate(self, topic, num_slides=5):
        """完整的生成流程"""
        print("=" * 70)
        print(f"🚀 SlideCraft AI 启动!")
        print(f"   主题: {topic}")
        print(f"   页数: {num_slides}")
        print(f"   模型: {self.model}")
        print("=" * 70)

        try:
            # 1. 生成大纲
            outline = self.generate_outline(topic, num_slides)

            # 2. 为每页生成内容
            print(f"\n📝 正在生成各页内容...")
            contents = []
            for slide_info in outline["slides"]:
                content = self.generate_slide_content(slide_info, topic)
                contents.append(content)

            # 3. 创建PPT
            output_path = self.create_ppt(outline, contents)

            print("\n" + "=" * 70)
            print(f"🎉 全部完成!")
            print(f"   生成的PPT: {output_path}")
            print("=" * 70)

            return output_path

        except Exception as e:
            print(f"\n❌ 生成失败: {str(e)}")
            import traceback
            traceback.print_exc()
            raise


def main():
    """主函数"""
    # 创建输出目录
    os.makedirs("output", exist_ok=True)

    # 创建原型实例
    crafter = SlideCraftPrototype()

    # 测试主题
    topics = [
        "人工智能在教育领域的应用",
        # 可以添加更多测试主题
    ]

    for topic in topics:
        try:
            crafter.generate(topic, num_slides=5)
            print("\n✅ 测试成功!\n")
        except Exception as e:
            print(f"\n❌ 测试失败: {str(e)}\n")


if __name__ == "__main__":
    main()
